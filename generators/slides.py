# -*- coding: utf-8 -*-
"""BOOK 6 — Teaching Slides: one PPTX deck per teaching session."""
import os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from curriculum import load_units, load_reviews
from curriculum.audio_links import lesson_audio, link_set

W, H = Inches(13.333), Inches(7.5)

# --- palette --------------------------------------------------------------
C = {
    "navy":  RGBColor(0x14, 0x2C, 0x4F),
    "blue":  RGBColor(0x1E, 0x6F, 0xB8),
    "sky":   RGBColor(0xE7, 0xF1, 0xFA),
    "teal":  RGBColor(0x0E, 0x7C, 0x66),
    "mint":  RGBColor(0xE3, 0xF3, 0xEE),
    "orange":RGBColor(0xE0, 0x7A, 0x1F),
    "cream": RGBColor(0xFD, 0xF2, 0xE0),
    "purple":RGBColor(0x6B, 0x3F, 0xA0),
    "lilac": RGBColor(0xF1, 0xEA, 0xF9),
    "red":   RGBColor(0xC0, 0x39, 0x2B),
    "rose":  RGBColor(0xFC, 0xEA, 0xE8),
    "green": RGBColor(0x2E, 0x7D, 0x32),
    "leaf":  RGBColor(0xEA, 0xF5, 0xEA),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "ink":   RGBColor(0x1B, 0x1B, 0x1B),
    "grey":  RGBColor(0x6B, 0x72, 0x80),
    "paper": RGBColor(0xFA, 0xFB, 0xFC),
}

SECTION = {   # section -> (accent, tint, icon)
    "warmup":   ("orange", "cream", "🔥"),
    "objectives": ("blue", "sky", "🎯"),
    "vocab":    ("orange", "cream", "📕"),
    "grammar":  ("teal", "mint", "🔧"),
    "pron":     ("purple", "lilac", "🔊"),
    "listening":("blue", "sky", "🎧"),
    "reading":  ("green", "leaf", "📖"),
    "speaking": ("red", "rose", "💬"),
    "writing":  ("navy", "sky", "✍️"),
    "practice": ("teal", "mint", "🧩"),
    "communication": ("red", "rose", "🌍"),
    "review":   ("green", "leaf", "🔁"),
    "homework": ("navy", "sky", "🏠"),
    "project":  ("purple", "lilac", "🛠"),
    "warn":     ("red", "rose", "⚠"),
}


# --- primitives -----------------------------------------------------------
def _txbox(slide, x, y, w, h, text, size=20, bold=False, color="ink", align=PP_ALIGN.LEFT,
           anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri", spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        r = p.add_run(); r.text = str(ln)
        f = r.font
        f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = C[color] if isinstance(color, str) else color
        f.name = font
    return tb


def _rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.06):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = C[fill] if isinstance(fill, str) else fill
    if line:
        sh.line.color.rgb = C[line] if isinstance(line, str) else line
        sh.line.width = Pt(1.5)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sh.adjustments[0] = adj
    except Exception:
        pass
    sh.text_frame.word_wrap = True
    return sh


def _blank(prs, bg="paper"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgsh = _rect(s, 0, 0, W, H, bg, shape=MSO_SHAPE.RECTANGLE)
    bgsh.line.fill.background()
    return s


def _header(slide, section, title, footer=""):
    accent, tint, icon = SECTION.get(section, ("blue", "sky", "•"))
    _rect(slide, 0, 0, W, Inches(1.15), accent, shape=MSO_SHAPE.RECTANGLE)
    _rect(slide, 0, Inches(1.15), W, Inches(0.06), "navy", shape=MSO_SHAPE.RECTANGLE)
    _txbox(slide, Inches(0.45), Inches(0.16), Inches(1.0), Inches(0.85), icon, size=34,
           color="white", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tsize = 28 if len(str(title)) <= 46 else (23 if len(str(title)) <= 66 else 20)
    _txbox(slide, Inches(1.45), Inches(0.16), Inches(9.8), Inches(0.85), title, size=tsize, bold=True,
           color="white", anchor=MSO_ANCHOR.MIDDLE)
    if footer:
        _txbox(slide, Inches(11.4), Inches(0.3), Inches(1.6), Inches(0.6), footer, size=12,
               color="white", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return tint, accent


def _wrapped(items, per=5):
    return [items[i:i + per] for i in range(0, len(items), per)] or [[]]


def _clean(s, limit=190):
    s = re.sub(r"\s+", " ", str(s)).strip()
    return s if len(s) <= limit else s[:limit - 1] + "…"


# --- slide builders -------------------------------------------------------
def title_slide(prs, L):
    s = _blank(prs, "navy")
    _rect(s, 0, 0, W, H, "navy", shape=MSO_SHAPE.RECTANGLE)
    _rect(s, Inches(-1.2), Inches(4.8), Inches(7.5), Inches(4.5), "blue", shape=MSO_SHAPE.OVAL)
    _rect(s, Inches(9.6), Inches(-2.0), Inches(6.0), Inches(5.0), "teal", shape=MSO_SHAPE.OVAL)
    _txbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6),
           f"UNIT {L.unit}  ·  LESSON {L.number}  ·  PERIOD {L.period}", size=20, bold=True, color="white")
    _txbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6), L.lesson_type.upper(),
           size=30, bold=True, color="sky")
    _txbox(s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(2.2), L.title, size=44, bold=True, color="white")
    _txbox(s, Inches(0.9), Inches(5.8), Inches(11.5), Inches(0.8),
           "English 7  ·  Let's start! 👋", size=20, color="sky")
    return s


def objectives_slide(prs, L):
    s = _blank(prs)
    tint, accent = _header(s, "objectives", "Today we will learn to…", L.code)
    y = Inches(1.6)
    for i, o in enumerate(L.objectives[:5]):
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(0.95), tint)
        _rect(s, Inches(0.7), y, Inches(0.14), Inches(0.95), accent, shape=MSO_SHAPE.RECTANGLE)
        _txbox(s, Inches(1.05), y + Inches(0.05), Inches(1.0), Inches(0.85), f"{i+1}", size=26,
               bold=True, color=accent, anchor=MSO_ANCHOR.MIDDLE)
        _txbox(s, Inches(1.9), y + Inches(0.05), Inches(10.4), Inches(0.85), _clean(o, 140), size=19,
               anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.08)
    return s


def warmup_slide(prs, L):
    st = L.procedure[0] if L.procedure else None
    if not st:
        return
    s = _blank(prs)
    tint, accent = _header(s, "warmup", f"Warm-up: {st.name.replace('Warm-up: ','')}  ({st.minutes}')", L.code)
    _rect(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.6), tint)
    y = Inches(1.9)
    for t in st.teacher[:4]:
        _txbox(s, Inches(1.1), y, Inches(11.1), Inches(0.9), "• " + _clean(t, 170), size=19)
        y += Inches(0.95)
    _txbox(s, Inches(1.1), Inches(6.4), Inches(11.1), Inches(0.7),
           f"👉 You: {_clean(st.students, 150)}", size=18, italic=True, color="grey")
    return s


def vocab_slides(prs, L):
    out = []
    for chunk in _wrapped(L.vocab, 3):
        if not chunk:
            break
        s = _blank(prs)
        tint, accent = _header(s, "vocab", "New words", L.code)
        y = Inches(1.55)
        for w in chunk:
            _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.62), tint)
            _txbox(s, Inches(0.95), y + Inches(0.12), Inches(4.6), Inches(0.7), w.word, size=30,
                   bold=True, color=accent)
            _txbox(s, Inches(0.95), y + Inches(0.82), Inches(4.6), Inches(0.6),
                   f"{w.ipa}   ({w.pos})", size=17, color="grey")
            _txbox(s, Inches(5.7), y + Inches(0.12), Inches(6.7), Inches(0.66), _clean(w.vn, 42),
                   size=19, bold=True, color="navy")
            _txbox(s, Inches(5.7), y + Inches(0.85), Inches(6.7), Inches(0.6),
                   "e.g. " + _clean(w.example, 90), size=15, italic=True, color="ink")
            y += Inches(1.75)
        out.append(s)
    if L.phrases:
        s = _blank(prs)
        tint, accent = _header(s, "vocab", "Useful phrases — learn them as one block!", L.code)
        y = Inches(1.7)
        for ph in L.phrases[:7]:
            _rect(s, Inches(0.9), y, Inches(11.5), Inches(0.68), "white", line=accent)
            _txbox(s, Inches(1.2), y + Inches(0.06), Inches(11.0), Inches(0.6), "🔹 " + _clean(ph, 110),
                   size=20, anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(0.78)
        out.append(s)
    return out


def grammar_slides(prs, L):
    out = []
    g = L.grammar
    if not g:
        return out
    s = _blank(prs)
    tint, accent = _header(s, "grammar", f"Grammar: {_clean(g.point, 60)}", L.code)
    y = Inches(1.6)
    for u in g.use[:4]:
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(1.0), tint)
        _txbox(s, Inches(1.05), y + Inches(0.08), Inches(11.3), Inches(0.85), "▸ " + _clean(u, 160),
               size=19, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.12)
    out.append(s)

    if g.form:
        s = _blank(prs)
        _header(s, "grammar", "The form", L.code)
        rows, cols = len(g.form), max(len(r) for r in g.form)
        top = Inches(1.6)
        height = min(Inches(5.3), Inches(0.62) * rows)
        shape = s.shapes.add_table(rows, cols, Inches(0.6), top, Inches(12.1), height)
        tbl = shape.table
        for i, row in enumerate(g.form):
            for j in range(cols):
                cell = tbl.cell(i, j)
                cell.text = _clean(row[j], 95) if j < len(row) else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(16 if rows > 4 else 18)
                        r.font.name = "Calibri"
                        r.font.bold = (i == 0)
                        r.font.color.rgb = C["white"] if i == 0 else C["ink"]
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["teal"] if i == 0 else (
                    C["mint"] if i % 2 else C["white"])
        out.append(s)

    if g.examples:
        s = _blank(prs)
        tint, accent = _header(s, "grammar", "Look at the examples", L.code)
        y = Inches(1.7)
        for e in g.examples[:5]:
            _rect(s, Inches(0.8), y, Inches(11.7), Inches(0.85), "white", line=accent)
            _txbox(s, Inches(1.15), y + Inches(0.08), Inches(11.2), Inches(0.7), "✅ " + _clean(e, 120),
                   size=20, anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(0.98)
        out.append(s)

    if g.pitfall:
        s = _blank(prs)
        tint, accent = _header(s, "warn", "Be careful! A very common mistake", L.code)
        _rect(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.6), "rose")
        _txbox(s, Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.2), _clean(g.pitfall, 330), size=21)
        if g.examples:
            _txbox(s, Inches(1.2), Inches(4.8), Inches(11.0), Inches(1.4),
                   "👍 Say it like this:  " + _clean(g.examples[0], 110), size=22, bold=True, color="green")
        out.append(s)
    return out


def pron_slide(prs, L):
    if not L.pron:
        return []
    p = L.pron
    s = _blank(prs)
    tint, accent = _header(s, "pron", f"Pronunciation: {_clean(p.focus, 55)}", L.code)
    _rect(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.35), tint)
    _txbox(s, Inches(1.05), Inches(1.65), Inches(11.3), Inches(1.2), _clean(p.tip, 260), size=18,
           anchor=MSO_ANCHOR.MIDDLE)
    y = Inches(3.15)
    for it in p.items[:4]:
        _rect(s, Inches(0.8), y, Inches(11.7), Inches(0.68), "white", line=accent)
        _txbox(s, Inches(1.1), y + Inches(0.05), Inches(11.2), Inches(0.6), "🔊 " + _clean(it, 110),
               size=19, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.78)
    out = [s]
    if p.drill:
        s2 = _blank(prs)
        tint, accent = _header(s2, "pron", "Repeat after me — three times each!", L.code)
        y = Inches(1.9)
        for d in p.drill[:4]:
            _rect(s2, Inches(0.8), y, Inches(11.7), Inches(1.0), tint)
            _txbox(s2, Inches(1.15), y + Inches(0.1), Inches(11.2), Inches(0.85), "🗣 " + _clean(d, 130),
                   size=21, anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(1.15)
        out.append(s2)
    return out


def _task_slides(prs, e, section, heading):
    """Turn one exercise into 1–3 slides (instruction + items)."""
    out = []
    chunks = _wrapped(e.items, 5)
    for k, chunk in enumerate(chunks):
        s = _blank(prs)
        tint, accent = _header(s, section, f"{heading}" + (f"  ({k+1}/{len(chunks)})" if len(chunks) > 1 else ""),
                               e.ref)
        _rect(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.85), tint)
        _txbox(s, Inches(1.0), Inches(1.58), Inches(11.3), Inches(0.75),
               _clean(e.instruction, 170), size=19, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        y = Inches(2.6)
        if e.wordbank and k == 0:
            _rect(s, Inches(0.9), y, Inches(11.5), Inches(0.7), "cream", line="orange")
            _txbox(s, Inches(1.2), y + Inches(0.05), Inches(11.0), Inches(0.6),
                   "   ".join(e.wordbank), size=18, bold=True, color="orange", anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(0.9)
        for it in chunk:
            _txbox(s, Inches(1.1), y, Inches(11.2), Inches(0.7), _clean(it, 150), size=19)
            y += Inches(0.72)
        out.append(s)
    if not e.items:
        s = _blank(prs)
        tint, accent = _header(s, section, heading, e.ref)
        _rect(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.4), tint)
        _txbox(s, Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.0), _clean(e.instruction, 300), size=22,
               anchor=MSO_ANCHOR.MIDDLE)
        out.append(s)
    return out



def _button_row(s, items, y, bw, bh, fill, line, colour, fontsize, prefix, single):
    x = Inches(0.8)
    gap = Inches(0.12)
    for label, target in items:
        box = _rect(s, x, y, bw, bh, fill, line=line)
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = single if len(items) == 1 else f"{prefix}{label}"
        run.font.size = Pt(fontsize)
        run.font.bold = True
        run.font.color.rgb = C[colour]
        run.hyperlink.address = target
        x = x + bw + gap


def _button_rows(local, online):
    """How many rows the two tiers will need (max 3 per row)."""
    return ((len(local) + 2) // 3 if local else 0,
            (len(online) + 2) // 3 if online else 0)


def _play_buttons(s, local, online, top, bh):
    """Two routes to the recording, as two rows of clickable buttons.

    Row 1 (blue)  the local MP3 — one button per part, because a Looking Back
                  lesson replaying a four-part recording needs five of them and
                  dropping one would silently lose a required part.
    Row 2 (teal)  the publisher's page — one button per DISTINCT page, since a
                  multipart recording is published as several conversations on
                  a single page. Four identical buttons would be noise.

    Both rows are always shown when both exist: a deck sent to a colleague
    travels without audio/, and the online button is what still works for them.

    Geometry is computed so the credit line below never leaves the slide.
    """
    if not local and not online:
        return top
    rows_l, rows_o = _button_rows(local, online)
    gap_y = Inches(0.07)
    usable = Inches(11.7)
    gap = Inches(0.12)

    y = top
    for r in range(rows_l):
        chunk = local[r * 3:(r + 1) * 3]
        bw = int((usable - gap * (len(chunk) - 1)) / len(chunk))
        _button_row(s, chunk, y, bw, bh, "sky", "blue", "navy",
                    15 if len(chunk) == 1 else 13, "▶ ", "▶ Play audio")
        y += bh + gap_y
    for r in range(rows_o):
        chunk = online[r * 3:(r + 1) * 3]
        bw = int((usable - gap * (len(chunk) - 1)) / len(chunk))
        _button_row(s, chunk, y, bw, bh, "mint", "teal", "teal",
                    14 if len(chunk) == 1 else 12, "🌐 ", "🌐 Listen online")
        y += bh + gap_y
    return y - gap_y      # the loop adds a trailing gap after the last row


def listening_slides(prs, L):
    if not L.listening:
        return []
    a = L.listening
    s = _blank(prs)
    tint, accent = _header(s, "listening", f"Listening: {_clean(a.title, 50)}", L.code)
    _rect(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.0), tint)
    _txbox(s, Inches(1.2), Inches(2.1), Inches(11.0), Inches(1.7), _clean(a.context, 260), size=21,
           anchor=MSO_ANCHOR.MIDDLE)

    # Lay the lower half out from the top. A Looking Back lesson replaying a
    # four-part recording needs five local buttons plus two online ones, which
    # is three rows; the "Listen 1/2/3" block gives up the height for them
    # rather than being overlapped by them.
    lo, on = link_set(lesson_audio(L.code), 3)
    rows_l, rows_o = _button_rows(lo, on)
    rows = rows_l + rows_o
    bh = Inches(0.52) if rows >= 3 else Inches(0.56)
    buttons_h = rows * bh + max(rows - 1, 0) * Inches(0.07)
    credit_h = Inches(0.42)
    inst_top = Inches(4.05)
    inst_h = Inches(7.45) - inst_top - Inches(0.12) - buttons_h - Inches(0.10) - credit_h
    inst_h = max(min(inst_h, Inches(1.6)), Inches(0.95))
    _txbox(s, Inches(0.8), inst_top, Inches(11.7), inst_h,
           ["🎧  Listen 1: just listen — do not write.",
            "🎧  Listen 2: write your answers.",
            "🎧  Listen 3: check with your partner."],
           size=20 if inst_h >= Inches(1.35) else 16, color="navy")

    below = _play_buttons(s, lo, on, inst_top + inst_h + Inches(0.12), bh)
    credit = " · ".join(x for x in (a.source, a.duration, a.speakers) if x)
    if credit:
        _txbox(s, Inches(0.8), below + Inches(0.10), Inches(11.7), credit_h,
               [_clean(credit, 170)], size=10, color="grey")
    out = [s]
    for t in a.tasks:
        out += _task_slides(prs, t, "listening", "Listening task")
    return out


def reading_slides(prs, L):
    if not L.reading:
        return []
    r = L.reading
    out = []
    s = _blank(prs)
    tint, accent = _header(s, "reading", f"Reading: {_clean(r.title, 50)}", L.code)
    body = " ".join(r.body)
    _rect(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.3), tint)
    _txbox(s, Inches(1.05), Inches(1.7), Inches(11.3), Inches(4.9), _clean(body, 900), size=15,
           spacing=1.15)
    out.append(s)
    for t in r.tasks:
        out += _task_slides(prs, t, "reading", "Reading task")
    return out


def practice_slides(prs, L):
    out = []
    for e in L.guided:
        out += _task_slides(prs, e, "practice", "Guided practice")
    for e in L.independent:
        out += _task_slides(prs, e, "practice", "Your turn — independent practice")
    return out


def speaking_slides(prs, L):
    out = []
    for e in L.speaking:
        out += _task_slides(prs, e, "speaking", "Speaking")
    if L.communication:
        c = L.communication
        s = _blank(prs)
        tint, accent = _header(s, "communication", f"Real-life English: {_clean(c.get('function',''), 55)}",
                               L.code)
        y = Inches(1.6)
        for ph in c.get("phrases", [])[:6]:
            _rect(s, Inches(0.9), y, Inches(11.5), Inches(0.72), "white", line=accent)
            _txbox(s, Inches(1.2), y + Inches(0.05), Inches(11.0), Inches(0.62), "💬 " + _clean(ph, 100),
                   size=20, anchor=MSO_ANCHOR.MIDDLE)
            y += Inches(0.82)
        out.append(s)
        if c.get("roleplay"):
            s2 = _blank(prs)
            tint, accent = _header(s2, "communication", "Role play — in pairs", L.code)
            _rect(s2, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.8), tint)
            _txbox(s2, Inches(1.2), Inches(2.1), Inches(11.0), Inches(2.4), _clean(c["roleplay"], 330),
                   size=21, anchor=MSO_ANCHOR.MIDDLE)
            _txbox(s2, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.4),
                   ["⏱ 2 minutes to prepare  ·  2 minutes to perform  ·  then swap roles",
                    "Do not read! Look at your partner."], size=19, color="navy", bold=True)
            out.append(s2)
    return out


def writing_slides(prs, L):
    out = []
    for e in L.writing:
        out += _task_slides(prs, e, "writing", "Writing")
    return out


def review_slide(prs, L):
    s = _blank(prs)
    tint, accent = _header(s, "review", "Let's remember today's English", L.code)
    y = Inches(1.75)
    for r in (L.review or [o for o in L.objectives])[:5]:
        _rect(s, Inches(0.8), y, Inches(11.7), Inches(0.92), tint)
        _txbox(s, Inches(1.15), y + Inches(0.08), Inches(11.2), Inches(0.78), "✓ " + _clean(r, 150),
               size=20, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.05)
    return s


def homework_slide(prs, L):
    s = _blank(prs)
    tint, accent = _header(s, "homework", "Homework — write it in your notebook", L.code)
    y = Inches(1.7)
    for h in L.homework[:5]:
        _rect(s, Inches(0.8), y, Inches(11.7), Inches(0.95), tint)
        _txbox(s, Inches(1.1), y + Inches(0.06), Inches(2.0), Inches(0.85), h.ref, size=19, bold=True,
               color=accent, anchor=MSO_ANCHOR.MIDDLE)
        _txbox(s, Inches(3.0), y + Inches(0.06), Inches(9.3), Inches(0.85),
               _clean(f"{h.title} — {h.instruction}", 150), size=17, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.06)
    _txbox(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.7),
           "Homework Book, section " + L.code + "   ·   See you next lesson! 👋", size=18,
           color="grey", align=PP_ALIGN.CENTER)
    return s


def build_deck(L, path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    title_slide(prs, L)
    objectives_slide(prs, L)
    warmup_slide(prs, L)
    vocab_slides(prs, L)
    grammar_slides(prs, L)
    pron_slide(prs, L)
    listening_slides(prs, L)
    reading_slides(prs, L)
    practice_slides(prs, L)
    speaking_slides(prs, L)
    writing_slides(prs, L)
    review_slide(prs, L)
    homework_slide(prs, L)
    prs.save(path)
    return path


def _slug(s):
    return re.sub(r"[^A-Za-z0-9]+", "_", s).strip("_")[:48]


def build(root="output/slides"):
    made = []
    groups = [(f"Unit{u.number:02d}_{_slug(u.title)}", u.lessons) for u in load_units()]
    groups += [(f"Review{r.number:02d}_{_slug(r.title)}", r.lessons) for r in load_reviews()]
    for folder, lessons in groups:
        d = os.path.join(root, folder)
        os.makedirs(d, exist_ok=True)
        for L in lessons:
            fn = f"{L.code}_{_slug(L.lesson_type)}.pptx"
            made.append(build_deck(L, os.path.join(d, fn)))
    return made
