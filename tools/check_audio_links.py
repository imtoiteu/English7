# -*- coding: utf-8 -*-
"""Verify every audio hyperlink in every generated DOCX and PPTX.

Answers, in order, the questions that actually matter:

  * does every listening activity carry a clickable link?
  * does every link resolve to a file that exists?
  * does any link point at the WRONG recording?
  * did linking change any existing content?

The last one is the reason a text baseline is compared rather than trusted: a
hyperlink is added by editing the document XML, and the cheapest way to be sure
nothing else moved is to diff the extracted text against a snapshot.

    python3 tools/check_audio_links.py
    python3 tools/check_audio_links.py --baseline path/to/text.json
"""
import os, re, sys, json, glob, zipfile, argparse, difflib

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)

from curriculum import all_lessons
from curriculum.audio_links import lesson_audio, diagnostic_audio, link_set, exists
from curriculum.diagnostic import ALL_PAPERS

FAILS, WARNS, OKS = [], [], []
def ok(m):   OKS.append(m)
def warn(m): WARNS.append(m)
def fail(m): FAILS.append(m)

DOCX = sorted(glob.glob(os.path.join(ROOT, "output", "*.docx")))
PPTX = sorted(glob.glob(os.path.join(ROOT, "output", "slides", "*", "*.pptx")))


# --------------------------------------------------------------------------
def external_links(path):
    """Every external relationship target in an OOXML package, with its part."""
    out = []
    with zipfile.ZipFile(path) as z:
        for n in z.namelist():
            if not n.endswith(".rels"):
                continue
            xml = z.read(n).decode("utf-8", "replace")
            for m in re.finditer(r'<Relationship\b[^>]*>', xml):
                tag = m.group(0)
                if 'TargetMode="External"' not in tag:
                    continue
                t = re.search(r'Target="([^"]+)"', tag)
                if t:
                    out.append((n, t.group(1)))
    return out


def audio_links_of(path):
    return [(part, t) for part, t in external_links(path) if t.lower().endswith(".mp3")]


def web_links_of(path):
    return [(part, t) for part, t in external_links(path)
            if t.startswith(("http://", "https://"))]


def expected_pages():
    """lesson/task code -> set of official source pages it should link."""
    exp = {}
    for L in all_lessons():
        if L.code.startswith("D0") or not L.listening:
            continue
        exp[L.code] = {u for _, u in link_set(lesson_audio(L.code), 1)[1]}
    for P in ALL_PAPERS:
        for s in P.sections:
            if s.strand != "listening":
                continue
            for t in s.tasks:
                exp[t.code] = {u for _, u in link_set(diagnostic_audio(t.audio_key), 1)[1]}
    return exp


def check_online():
    """Both routes must exist, and the online one must be the exact page."""
    from curriculum.audio_sources import AUDIO
    from curriculum.audio_diagnostic import DIAG_AUDIO
    legit = {a.source_page for a in AUDIO.values() if a.source_page}
    legit |= {a.source_page for a in DIAG_AUDIO.values() if a.source_page}

    generic = re.compile(r"^https?://(www\.)?(elllo\.org|learningenglish\.voanews\.com)/?$")
    total = 0
    for p in DOCX + PPTX:
        for _, u in web_links_of(p):
            total += 1
            if generic.match(u):
                fail(f"{os.path.relpath(p, ROOT)} links a generic homepage: {u}")
            elif u not in legit:
                fail(f"{os.path.relpath(p, ROOT)} links {u}, which is not a recorded source page")
    if not [f for f in FAILS if "generic homepage" in f or "not a recorded source" in f]:
        ok(f"{total} online links, every one an exact recorded source page (no homepages)")

    # each deck's online set must be exactly its lesson's
    exp = expected_pages()
    wrong = 0
    for p in PPTX:
        code = os.path.basename(p).split("_")[0]
        want = exp.get(code)
        if want is None:
            continue
        got = {u for _, u in web_links_of(p)}
        if got != want:
            wrong += 1
            fail(f"{os.path.relpath(p, ROOT)}: online links {sorted(got)} but {code} "
                 f"needs {sorted(want)}")
    if not wrong:
        ok(f"every one of the {len(PPTX)} decks links exactly its lesson's official page(s)")

    # both routes present wherever both are possible
    for p in DOCX + PPTX:
        n_local, n_web = len(audio_links_of(p)), len(web_links_of(p))
        base = os.path.relpath(p, ROOT)
        if n_local and not n_web:
            fail(f"{base}: {n_local} local links but no online fallback")
    if not [f for f in FAILS if "no online fallback" in f]:
        ok("every document that offers a local MP3 also offers the online source")
    return total


def check_diagnostic_pages():
    """A-L1…C-L3 must carry their exact official VOA pages."""
    from curriculum.audio_diagnostic import DIAG_AUDIO
    from curriculum.diagnostic import PAPERS
    need = {}
    for P in PAPERS.values():
        for s in P.sections:
            if s.strand != "listening":
                continue
            for t in s.tasks:
                need[t.code] = DIAG_AUDIO[t.audio_key].source_page
    for label, name in (("Book 6", "06_Diagnostic_and_Adaptive_System.docx"),
                        ("Book 7", "07_Diagnostic_Test_Papers.docx")):
        p = os.path.join(ROOT, "output", name)
        if not os.path.exists(p):
            continue
        got = {u for _, u in web_links_of(p)}
        gap = {c: u for c, u in need.items() if u not in got}
        if gap:
            fail(f"{label}: diagnostic tasks missing their official page: {sorted(gap)}")
        else:
            ok(f"{label}: all 8 diagnostic tasks carry their exact official VOA page")


def resolves(doc_path, target):
    """Does a relative link resolve to a real file from the document's location?"""
    if target.startswith(("http://", "https://")):
        return None                      # a web fallback; existence not checked offline
    return os.path.isfile(os.path.normpath(os.path.join(os.path.dirname(doc_path), target)))


# --------------------------------------------------------------------------
def expected_map():
    """lesson/task code -> set of expected mp3 basenames."""
    exp = {}
    for L in all_lessons():
        if L.code.startswith("D0"):
            continue
        if L.listening:
            exp[L.code] = {fn for _, fn, _ in lesson_audio(L.code)}
    for P in ALL_PAPERS:
        for s in P.sections:
            if s.strand != "listening":
                continue
            for t in s.tasks:
                exp[t.code] = {fn for _, fn, _ in diagnostic_audio(t.audio_key)}
    return exp


def check_inventory():
    exp = expected_map()
    missing = [(k, f) for k, fs in exp.items() for f in fs if not exists(f)]
    if missing:
        for k, f in missing:
            fail(f"{k}: expected recording {f} is not in audio/")
    else:
        n = len({f for fs in exp.values() for f in fs})
        ok(f"inventory: {len(exp)} listening activities → {n} distinct MP3s, all present")
    # ambiguity: two different owners must never resolve to the same single file
    # unless that is a documented replay (Looking Back lessons)
    return exp


def check_resolution():
    bad = miss = total = 0
    for p in DOCX + PPTX:
        for part, t in audio_links_of(p):
            total += 1
            r = resolves(p, t)
            if r is False:
                bad += 1
                fail(f"{os.path.relpath(p, ROOT)} → {t} does not resolve to a file")
            elif r is None:
                miss += 1
    if not bad:
        ok(f"{total} audio hyperlinks across {len(DOCX)} DOCX and {len(PPTX)} PPTX, "
           f"every relative link resolves to a real MP3 on disk")
    if miss:
        warn(f"{miss} links are web fallbacks (the local MP3 was absent at build time)")
    return total


def check_correctness():
    """No link may point at a recording that belongs to a different lesson."""
    exp = expected_map()
    wrong = 0
    # slides: one deck per lesson, so the deck's links must be exactly that lesson's
    for p in PPTX:
        code = os.path.basename(p).split("_")[0]
        want = exp.get(code)
        if want is None:
            continue
        got = {os.path.basename(t) for _, t in audio_links_of(p)}
        if got != want:
            wrong += 1
            fail(f"{os.path.relpath(p, ROOT)}: links {sorted(got)} but {code} needs {sorted(want)}")
    if not wrong:
        ok(f"every one of the {len(PPTX)} decks links exactly the recordings its lesson uses")

    # books: the union of their links must be a subset of the whole inventory.
    # BRIDGE and EXT belong to no paper but are legitimately played by the
    # teacher (trigger T2 slow input, extension E5), so Book 6 links them too.
    from curriculum.audio_diagnostic import DIAG_FILES
    universe = {f for fs in exp.values() for f in fs}
    universe |= {os.path.basename(DIAG_FILES[k]) for k in ("BRIDGE", "EXT")}
    for p in DOCX:
        got = {os.path.basename(t) for _, t in audio_links_of(p)}
        stray = got - universe
        if stray:
            fail(f"{os.path.basename(p)} links recordings no activity uses: {sorted(stray)}")
    if not [f for f in FAILS if "no activity uses" in f]:
        ok("no book links a recording that belongs to no listening activity")


def check_coverage():
    """Every listening activity that is rendered must have got a link."""
    exp = expected_map()
    # slides
    decks = {os.path.basename(p).split("_")[0] for p in PPTX}
    lesson_codes = {c for c in exp if not c.startswith(("A-", "B-", "C-"))}
    uncovered = sorted(lesson_codes & decks - {os.path.basename(p).split("_")[0]
                                               for p in PPTX if audio_links_of(p)})
    if uncovered:
        fail(f"decks with a listening activity but no play button: {uncovered}")
    else:
        ok(f"all {len(lesson_codes & decks)} lesson decks with listening carry a play button")

    # diagnostic tasks must appear in both diagnostic books
    d6 = os.path.join(ROOT, "output", "06_Diagnostic_and_Adaptive_System.docx")
    d7 = os.path.join(ROOT, "output", "07_Diagnostic_Test_Papers.docx")
    for label, p in (("Book 6", d6), ("Book 7", d7)):
        if not os.path.exists(p):
            warn(f"{label} not built")
            continue
        got = {os.path.basename(t) for _, t in audio_links_of(p)}
        need = set()
        for code in ("A-L1", "A-L2", "A-L3", "B-L1", "B-L2", "C-L1", "C-L2", "C-L3"):
            need |= exp.get(code, set())
        gap = need - got
        if gap:
            fail(f"{label}: diagnostic recordings with no link: {sorted(gap)}")
        else:
            ok(f"{label}: all 8 diagnostic listening tasks (A-L1…C-L3) are linked")


def check_books_have_links():
    for p in DOCX:
        n = len(audio_links_of(p))
        base = os.path.basename(p)
        if base.startswith(("03_", "04_")):
            if n:
                warn(f"{base}: {n} audio links, though it renders no listening activity")
            else:
                ok(f"{base}: no listening activity rendered, correctly no links")
        elif n == 0:
            fail(f"{base}: renders listening activities but has no audio links")
        else:
            ok(f"{base}: {n} audio hyperlinks")


def check_content_unchanged(baseline):
    if not baseline or not os.path.exists(baseline):
        warn("no text baseline supplied; content-preservation check skipped")
        return
    from docx import Document
    from pptx import Presentation
    snap = json.load(open(baseline))

    def docx_text(p):
        d = Document(p); out = [x.text for x in d.paragraphs]
        for t in d.tables:
            for r in t.rows:
                for c in r.cells:
                    out.append(c.text)
        return out

    def pptx_text(p):
        pr = Presentation(p); out = []
        for i, s in enumerate(pr.slides):
            for sh in s.shapes:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        out.append(f"{i}|{para.text}")
        return out

    # A table cell's text is one blob containing several lines, so filtering
    # whole entries would discard an entire cell just because a link was added
    # inside it. Split first, filter the added lines, then compare.
    # Link text added by this feature. 🌐 must be here as well as the words
    # "Listen online": when a lesson has several distinct source pages the
    # buttons are labelled per page ("🌐 U7L2"), not with the generic wording.
    LINKY = re.compile(r"▶|🌐|Listen online|To check a listening answer")
    def flatten(seq):
        out = []
        for entry in seq:
            for line in str(entry).split("\n"):
                line = line.strip()
                if line and not LINKY.search(line):
                    out.append(line)
        return out

    def declared(rel, b, a):
        """One content change is intentional and is declared here rather than
        waved through: Book 6's pre-flight check table used to truncate every
        filename to 34 characters, which made it useless for finding a file.
        It now prints the full name and carries a Play column. Any OTHER
        difference in that document is still a failure."""
        if not rel.endswith("06_Diagnostic_and_Adaptive_System.docx"):
            return False
        sm = difflib.SequenceMatcher(None, b, a, autojunk=False)
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                continue
            old, new = b[i1:i2], a[j1:j2]
            if old == ["Local file"] and new == ["Local file (full name)"]:
                continue
            if not old and new == ["On disk"]:
                continue
            if old == ["✔ on disk"] and new == ["✔"]:
                continue
            if (len(old) == len(new) == 1 and old[0].endswith("…")
                    and new[0].startswith(old[0][:-1]) and new[0].endswith(".mp3")):
                continue
            return False
        return True

    changed = []
    for rel, before in snap.items():
        p = os.path.join(ROOT, rel)
        if not os.path.exists(p):
            changed.append((rel, "file missing now"))
            continue
        after = docx_text(p) if rel.endswith(".docx") else pptx_text(p)
        b = flatten(before)
        a = flatten(after)
        if b != a and declared(rel, b, a):
            ok("Book 6's pre-flight table now prints full filenames and a Play column "
               "(declared change); nothing else in that document differs")
            continue
        if b != a:
            i = next((k for k in range(min(len(a), len(b))) if a[k] != b[k]), min(len(a), len(b)))
            changed.append((rel, f"{len(b)}→{len(a)} lines; first diff at {i}: "
                                 f"{b[i][:60]!r} → {a[i][:60]!r}" if i < min(len(a), len(b))
                                 else f"{len(b)}→{len(a)} lines"))
    if changed:
        for rel, why in changed[:8]:
            fail(f"content changed in {rel}: {why}")
        if len(changed) > 8:
            fail(f"…and {len(changed) - 8} more files with changed content")
    else:
        ok(f"content preserved: all {len(snap)} documents match the pre-link baseline "
           f"once the added link text is discounted")


def check_files_open():
    from docx import Document
    from pptx import Presentation
    bad = []
    for p in DOCX:
        try: Document(p)
        except Exception as e: bad.append((p, e))
    for p in PPTX:
        try: Presentation(p)
        except Exception as e: bad.append((p, e))
    if bad:
        for p, e in bad[:5]:
            fail(f"{os.path.relpath(p, ROOT)} will not open: {e}")
    else:
        ok(f"all {len(DOCX)} DOCX and {len(PPTX)} PPTX open cleanly")


def check_no_embedded_audio():
    big = []
    for p in DOCX + PPTX:
        with zipfile.ZipFile(p) as z:
            media = [n for n in z.namelist()
                     if n.startswith(("word/media/", "ppt/media/")) and n.endswith(".mp3")]
            if media:
                big.append((p, len(media)))
    if big:
        for p, n in big:
            fail(f"{os.path.relpath(p, ROOT)} embeds {n} MP3(s); links were requested instead")
    else:
        ok("no document embeds an MP3 — every reference is a lightweight hyperlink")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--baseline", help="text.json snapshot taken before linking")
    args = ap.parse_args()

    check_inventory()
    total = check_resolution()
    check_online()
    check_diagnostic_pages()
    check_correctness()
    check_coverage()
    check_books_have_links()
    check_no_embedded_audio()
    check_files_open()
    check_content_unchanged(args.baseline)

    print("\n".join("  ✓ " + m for m in OKS))
    if WARNS:
        print("\n".join("  ! " + m for m in WARNS))
    if FAILS:
        print("\n".join("  ✗ " + m for m in FAILS))
    print(f"\n  {len(OKS)} passed · {len(WARNS)} warnings · {len(FAILS)} failures")
    sys.exit(1 if FAILS else 0)
